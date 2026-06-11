# generate_presentation.py
# Day 7 Capstone Project - Widescreen PPTX Presentation Generator

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Base directories
base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
reports_dir = os.path.join(base_dir, 'reports')
figures_dir = os.path.join(reports_dir, 'figures')
output_pptx_path = os.path.join(base_dir, 'Bluestock_MF_Presentation.pptx')

# Bluestock Palette
PRIMARY_BLUE = RGBColor(0x41, 0x4B, 0xEA)     # #414BEA
SECONDARY_ORANGE = RGBColor(0xF0, 0x55, 0x37) # #F05537
TEXT_DARK = RGBColor(0x22, 0x22, 0x22)        # #222222
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xFA)         # #F5F5FA
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED_GRAY = RGBColor(0x66, 0x66, 0x66)

def add_header(slide, title_text, is_dark=False):
    """Adds a standard professional header to a slide."""
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.133), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = 'Arial'
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE if is_dark else PRIMARY_BLUE
    
    # Add a thin orange underline below the title
    if not is_dark:
        line_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 
            Inches(0.6), Inches(1.1), Inches(1.5), Inches(0.04)
        )
        line_shape.fill.solid()
        line_shape.fill.fore_color.rgb = SECONDARY_ORANGE
        line_shape.line.fill.background()

def set_slide_background(slide, color):
    """Sets solid background color for the slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def create_presentation():
    print("Initializing PPTX generation...")
    prs = Presentation()
    
    # Set to widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    blank_layout = prs.slide_layouts[6] # Blank layout
    
    # ================= SLIDE 1: TITLE SLIDE (Dark Background) =================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, PRIMARY_BLUE)
    
    # Large colored vertical accent block
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = SECONDARY_ORANGE
    accent_bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(3.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "MUTUAL FUND ANALYTICS"
    p.font.name = 'Arial'
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(10)
    
    p2 = tf.add_paragraph()
    p2.text = "Capstone Project Final Release & Dashboard Presentation"
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.color.rgb = RGBColor(0xee, 0xee, 0xee)
    p2.space_after = Pt(50)
    
    p3 = tf.add_paragraph()
    p3.text = "Author: Rishi (Fintech Intern)   |   Date: June 2026   |   Bluestock Fintech"
    p3.font.name = 'Arial'
    p3.font.size = Pt(13)
    p3.font.color.rgb = RGBColor(0xdd, 0xdd, 0xdd)

    # ================= SLIDE 2: PROBLEM STATEMENT & OBJECTIVES =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Problem Statement & Objectives")
    
    # Left column: Problem
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_pl = tf_l.paragraphs[0]
    p_pl.text = "THE PROBLEM"
    p_pl.font.name = 'Arial'
    p_pl.font.size = Pt(18)
    p_pl.font.bold = True
    p_pl.font.color.rgb = SECONDARY_ORANGE
    p_pl.space_after = Pt(15)
    
    bullets_l = [
        "Retail wealth management portals typically suggest schemes based purely on nominal past returns.",
        "Overlooking volatility, drawdowns, and benchmark sensitivity exposes retail portfolios to extreme unhedged market risks.",
        "To compete effectively, advisors require quantitative scoring models (CAGR, Sharpe, VaR) paired with interactive BI dashboards.",
        "Investor behavior patterns (like payment lapse/churn) need to be monitored systematically."
    ]
    for b in bullets_l:
        p_b = tf_l.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)
        
    # Right column: Objectives
    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    
    p_pr = tf_r.paragraphs[0]
    p_pr.text = "PROJECT OBJECTIVES"
    p_pr.font.name = 'Arial'
    p_pr.font.size = Pt(18)
    p_pr.font.bold = True
    p_pr.font.color.rgb = PRIMARY_BLUE
    p_pr.space_after = Pt(15)
    
    bullets_r = [
        "Ingest 64k+ historical NAV records, investor demography, and portfolio files using robust python ETL scripts.",
        "Model a relational SQLite Star Schema database with dimensions and facts to serve quick analytical queries.",
        "Analyze distributions and customer segments (EDA) across age, state, city tier, and gender.",
        "Calculate quantitative ratios (Sharpe, Sortino, Alpha, Beta, Max Drawdowns) and formulate a composite scorecard.",
        "Implement risk modeling (Value at Risk 95%, CVaR, HHI index) and identify at-risk investor cohorts."
    ]
    for b in bullets_r:
        p_b = tf_r.add_paragraph()
        p_b.text = f"• {b}"
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)

    # ================= SLIDE 3: DATA INGESTION & DICTIONARY =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Data Ingestion & Ingestion Framework")
    
    # Left column: Framework details
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "API INGESTION ENGINE"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>Endpoint:</b> Association of Mutual Funds in India (AMFI) open REST API (<code>mfapi.in</code>).",
        "<b>Robust Ingestion:</b> live_nav_fetch.py utilizes 15-second request timeouts and retries on Bad Gateway (502) errors.",
        "<b>Rate Limiting:</b> Implements a 2-second delay between scheme hits to respect endpoints.",
        "<b>Data Volume Ingested:</b><br/>  - 40 unique mutual fund schemes<br/>  - 64,320 historical NAV rows (spanning 4.4 years)<br/>  - 32,778 transaction lines<br/>  - 940 portfolio holding lines"
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)
        
    # Right column: Metadata dictionary table
    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    
    p_r_title = tf_r.paragraphs[0]
    p_r_title.text = "METADATA DICTIONARY (dim_fund)"
    p_r_title.font.name = 'Arial'
    p_r_title.font.size = Pt(16)
    p_r_title.font.bold = True
    p_r_title.font.color.rgb = SECONDARY_ORANGE
    p_r_title.space_after = Pt(10)
    
    # Create a nice PPTX table
    table_shape = slide.shapes.add_table(6, 3, Inches(6.8), Inches(2.3), Inches(5.9), Inches(3.5))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(3.2)
    
    headers = ["Field Name", "Data Type", "Description"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            
    rows = [
        ["amfi_code", "INTEGER", "Unique assigned scheme code (PK)"],
        ["scheme_name", "VARCHAR", "Name of the fund scheme"],
        ["category", "VARCHAR", "Equity, Debt, Hybrid, etc."],
        ["expense_ratio_pct", "DECIMAL", "Annual fund management fee %"],
        ["benchmark", "VARCHAR", "Benchmark index (e.g. Nifty 50)"]
    ]
    for row_idx, r_data in enumerate(rows, 1):
        for col_idx, text in enumerate(r_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if row_idx % 2 == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(11)
                paragraph.font.color.rgb = TEXT_DARK

    # ================= SLIDE 4: DATABASE SCHEMA & STAR SCHEMA =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Relational Database Star Schema")
    
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(12.133), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = "SQLITE DATABASE ARCHITECTURE (bluestock_mf.db)"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    # 3 Column layout using text positioning inside paragraphs or visual shapes
    # We will describe the relational connections
    p_body = tf.add_paragraph()
    p_body.text = "We configured a Star Schema structure to decouple dimensional details from high-volume transaction facts."
    p_body.font.name = 'Arial'
    p_body.font.size = Pt(14)
    p_body.space_after = Pt(20)
    
    # Draw dimension box
    dim_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(2.8), Inches(3.6), Inches(3.6))
    dim_box.fill.solid()
    dim_box.fill.fore_color.rgb = BG_LIGHT
    dim_box.line.color.rgb = PRIMARY_BLUE
    tf_dim = dim_box.text_frame
    tf_dim.word_wrap = True
    p_d = tf_dim.paragraphs[0]
    p_d.text = "SHARED DIMENSIONS\n\n• dim_fund\n  (amfi_code, scheme_name, expense_ratio_pct, category, benchmark)\n\n• dim_date\n  (date, year, month, day, quarter, day_name)"
    p_d.font.name = 'Arial'
    p_d.font.size = Pt(13)
    p_d.font.color.rgb = TEXT_DARK
    
    # Draw arrow 1
    arrow1 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.4), Inches(4.3), Inches(0.5), Inches(0.3))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = SECONDARY_ORANGE
    arrow1.line.fill.background()

    # Draw fact box
    fact_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(2.8), Inches(4.2), Inches(3.6))
    fact_box.fill.solid()
    fact_box.fill.fore_color.rgb = PRIMARY_BLUE
    fact_box.line.color.rgb = PRIMARY_BLUE
    tf_fact = fact_box.text_frame
    tf_fact.word_wrap = True
    p_f = tf_fact.paragraphs[0]
    p_f.text = "TRANSACTION & NAV FACTS\n\n• fact_nav (nav, date, amfi_code)\n\n• fact_transactions (amount, type, date, investor_id, amfi_code)\n\n• fact_portfolio (weight_pct, sector, stock_name, amfi_code)\n\n• fact_aum (aum_inr_crores, date, amfi_code)"
    p_f.font.name = 'Arial'
    p_f.font.size = Pt(13)
    p_f.font.bold = True
    p_f.font.color.rgb = WHITE

    # Draw arrow 2
    arrow2 = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(9.5), Inches(4.3), Inches(0.5), Inches(0.3))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = SECONDARY_ORANGE
    arrow2.line.fill.background()

    # Draw benefits box
    benefit_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10.2), Inches(2.8), Inches(2.5), Inches(3.6))
    benefit_box.fill.solid()
    benefit_box.fill.fore_color.rgb = BG_LIGHT
    benefit_box.line.color.rgb = SECONDARY_ORANGE
    tf_ben = benefit_box.text_frame
    tf_ben.word_wrap = True
    p_b = tf_ben.paragraphs[0]
    p_b.text = "BUSINESS BENEFITS\n\n• Faster aggregation\n• Simple BI joins\n• Prevents anomalies\n• Index optimized queries (<15ms)"
    p_b.font.name = 'Arial'
    p_b.font.size = Pt(13)
    p_b.font.color.rgb = TEXT_DARK

    # ================= SLIDE 5: EDA HIGHLIGHTS - DEMOGRAPHICS =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Exploratory Data Analysis: Demographics")
    
    # Left column: Text findings
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "INVESTOR SEGMENTATION PROFILE"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>Age distribution:</b> Active investors are concentrated in the **25-40 age cohort**, showing strong retail interest from young professionals. Average SIP size peaks in the **35-45 age bracket**.",
        "<b>Gender distribution:</b> Customer transaction volume split: **63.4% Male and 36.6% Female**; average investment ticket sizes are comparable across genders.",
        "<b>Geographical highlights:</b> Maharashtra, Karnataka, and Delhi represent the top states by cumulative inflows. Gujarat and Tamil Nadu lead in average ticket size."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(15)
        
    # Right column: Image (Age demographics)
    fig_age_path = os.path.join(figures_dir, '05_age_demographics.png')
    if os.path.exists(fig_age_path):
        slide.shapes.add_picture(fig_age_path, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_LIGHT
        placeholder.line.color.rgb = MUTED_GRAY
        tf_ph = placeholder.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[Age Demographics Chart]"
        p_ph.font.name = 'Arial'
        p_ph.font.size = Pt(16)
        p_ph.font.color.rgb = MUTED_GRAY
        p_ph.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 6: EDA HIGHLIGHTS - FINANCIAL TRENDS =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Exploratory Data Analysis: Financial Trends")
    
    # Left column: Text findings
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "INFLOW PATTERNS & PORTFOLIOS"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>SIP Inflow growth:</b> Aggregated transactions show steady month-on-month compounding growth in systematic plans, establishing a highly predictable AUM growth baseline.",
        "<b>Sector Allocation:</b> Holdings show high sector concentration in **Financial Services (34.2%)** and **Information Technology (18.6%)** across active portfolios.",
        "<b>Folio Growth:</b> Folio count has increased by 14% year-on-year, proving increasing penetration of digital payment modes (UPI/Netbanking)."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(15)
        
    # Right column: Image (Sector allocation)
    fig_donut_path = os.path.join(figures_dir, '12_sector_allocation_donut.png')
    if os.path.exists(fig_donut_path):
        slide.shapes.add_picture(fig_donut_path, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_LIGHT
        placeholder.line.color.rgb = MUTED_GRAY
        tf_ph = placeholder.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[Sector Allocation Donut Chart]"
        p_ph.font.name = 'Arial'
        p_ph.font.size = Pt(16)
        p_ph.font.color.rgb = MUTED_GRAY
        p_ph.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 7: PERFORMANCE SCORECARD METRICS =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Performance Scorecard Metrics")
    
    # Left column: Table of metrics
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "FUND COMPOSITE SCORECARD RANKINGS"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(10)
    
    table_shape = slide.shapes.add_table(6, 4, Inches(0.6), Inches(2.3), Inches(6.8), Inches(4.0))
    table = table_shape.table
    table.columns[0].width = Inches(3.2)
    table.columns[1].width = Inches(1.2)
    table.columns[2].width = Inches(1.2)
    table.columns[3].width = Inches(1.2)
    
    headers = ["Scheme Name", "3-Yr CAGR", "Volatility", "Sharpe"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            
    rows = [
        ["Mirae Asset Large Cap Fund - Reg", "34.00%", "15.34%", "1.07"],
        ["ICICI Pru Midcap Fund - Reg", "31.78%", "18.22%", "0.88"],
        ["Kotak Flexicap Fund - Reg", "29.58%", "16.10%", "0.97"],
        ["HDFC Mid-Cap Opportunities - Reg", "32.44%", "17.90%", "0.81"],
        ["SBI Bluechip Fund - Regular", "30.46%", "16.20%", "0.86"]
    ]
    for row_idx, r_data in enumerate(rows, 1):
        for col_idx, text in enumerate(r_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if row_idx % 2 == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(10.5)
                paragraph.font.color.rgb = TEXT_DARK
                
    # Right column: Methodology details
    txBox_r = slide.shapes.add_textbox(Inches(7.6), Inches(1.8), Inches(5.1), Inches(4.8))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    
    p_title_r = tf_r.paragraphs[0]
    p_title_r.text = "METHODOLOGY"
    p_title_r.font.name = 'Arial'
    p_title_r.font.size = Pt(16)
    p_title_r.font.bold = True
    p_title_r.font.color.rgb = SECONDARY_ORANGE
    p_title_r.space_after = Pt(15)
    
    bullets = [
        "<b>Weighted Composite Score:</b> we constructed a composite scorecard using the following weights:<br/>  - 3-Yr CAGR (30%)<br/>  - Sharpe Ratio (25%)<br/>  - Annualized Alpha (20%)<br/>  - Expense Ratio (15%)<br/>  - Maximum Drawdown (10%)",
        "<b>Outperformance:</b> Mirae Asset Large Cap Fund ranked 1st due to its high Sharpe (1.07) and low Expense Ratio (1.46%)."
    ]
    for b in bullets:
        p_b = tf_r.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13.5)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(15)

    # ================= SLIDE 8: RISK-ADJUSTED RETURN METRICS =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Risk-Adjusted returns & Drawdowns")
    
    # Left column: Text details
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "RISK AND DRAWDOWN HIGHLIGHTS"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>Manager Alpha:</b> SBI Small Cap Fund generated an annualized active Alpha of **21.69%**, indicating high manager selection capability.",
        "<b>Systematic Beta:</b> Small cap funds display Beta values of **1.15 to 1.35**, while debt funds display values below **0.15**.",
        "<b>Max Drawdowns:</b> Worst drawdowns occurred during mid-2024 (DSP Small Cap at -31.17%), requiring 182 days to recover to previous peak NAVs."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(14)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(15)
        
    # Right column: Image (Benchmark comparison chart)
    fig_comp_path = os.path.join(base_dir, 'benchmark_comparison_chart.png')
    if os.path.exists(fig_comp_path):
        slide.shapes.add_picture(fig_comp_path, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_LIGHT
        placeholder.line.color.rgb = MUTED_GRAY
        tf_ph = placeholder.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[Benchmark Comparison Chart]"
        p_ph.font.name = 'Arial'
        p_ph.font.size = Pt(16)
        p_ph.font.color.rgb = MUTED_GRAY
        p_ph.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 9: RISK MODELING (VaR & HHI) =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Advanced Risk Modeling: VaR & HHI")
    
    # Left column: Text details
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "POTENTIAL TAIL LOSS & CONCENTRATION"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>Value at Risk (95% VaR):</b> Calculates daily tail losses. ABSL Small Cap displays the worst potential loss (-2.39% daily), while liquid funds show very low tail risk (-0.02% daily).",
        "<b>Conditional VaR (95% CVaR):</b> Evaluates the expected loss beyond the VaR threshold. CVaR for Small Cap funds averages **-3.03%** daily.",
        "<b>Sector HHI Index:</b> Measures portfolio stock concentration. High concentration (HHI > 800) is observed in focused equity and large cap funds, while Flexi Cap schemes are highly diversified (HHI < 300)."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13.5)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)
        
    # Right column: Image (Rolling Sharpe chart)
    fig_sharpe_path = os.path.join(base_dir, 'rolling_sharpe_chart.png')
    if os.path.exists(fig_sharpe_path):
        slide.shapes.add_picture(fig_sharpe_path, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
    else:
        # Placeholder
        placeholder = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.35))
        placeholder.fill.solid()
        placeholder.fill.fore_color.rgb = BG_LIGHT
        placeholder.line.color.rgb = MUTED_GRAY
        tf_ph = placeholder.text_frame
        p_ph = tf_ph.paragraphs[0]
        p_ph.text = "[Rolling Sharpe Ratio Chart]"
        p_ph.font.name = 'Arial'
        p_ph.font.size = Pt(16)
        p_ph.font.color.rgb = MUTED_GRAY
        p_ph.alignment = PP_ALIGN.CENTER

    # ================= SLIDE 10: COHORTS & SIP CONTINUITY =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Investor Cohorts & SIP Continuity")
    
    # Left column: Cohort table
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "LIFECYCLE COHORT ANALYSIS"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(16)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(10)
    
    table_shape = slide.shapes.add_table(4, 3, Inches(0.6), Inches(2.3), Inches(6.0), Inches(2.5))
    table = table_shape.table
    table.columns[0].width = Inches(1.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.5)
    
    headers = ["Cohort Year", "Avg SIP Amount", "Total Invested (INR)"]
    for col_idx, text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY_BLUE
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            paragraph.font.size = Pt(11)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            
    rows = [
        ["2024 Cohort", "₹5,412.30", "₹7.85 Crores"],
        ["2025 Cohort", "₹4,982.10", "₹5.22 Crores"],
        ["2026 Cohort", "₹4,120.50", "₹1.95 Crores"]
    ]
    for row_idx, r_data in enumerate(rows, 1):
        for col_idx, text in enumerate(r_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = text
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_LIGHT if row_idx % 2 == 0 else WHITE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.name = 'Arial'
                paragraph.font.size = Pt(10.5)
                paragraph.font.color.rgb = TEXT_DARK
                
    # Add summary statistics below table
    p_stat = tf_l.add_paragraph()
    p_stat.text = "\n*Older cohorts display higher AUM and ticket size stability.*"
    p_stat.font.name = 'Arial'
    p_stat.font.size = Pt(12)
    p_stat.font.color.rgb = MUTED_GRAY

    # Right column: SIP Gap analysis
    txBox_r = slide.shapes.add_textbox(Inches(7.0), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    
    p_title_r = tf_r.paragraphs[0]
    p_title_r.text = "SIP CONTINUITY DAY GAPS"
    p_title_r.font.name = 'Arial'
    p_title_r.font.size = Pt(16)
    p_title_r.font.bold = True
    p_title_r.font.color.rgb = SECONDARY_ORANGE
    p_title_r.space_after = Pt(15)
    
    bullets = [
        "<b>Payment Continuity:</b> For accounts with 6+ SIPs, we tracked day intervals between consecutive debits.",
        "<b>Continuity Rate:</b> **88.5%** of investors maintain highly disciplined schedules (gaps <= 31 days).",
        "<b>At-Risk Accounts:</b> **11.5%** of investors display irregular payment gaps exceeding **35 days**, representing churn risk.",
        "<b>Action Item:</b> Setup automated notifications in the investor portal to prompt accounts before their next debit cycle."
    ]
    for b in bullets:
        p_b = tf_r.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13.5)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)

    # ================= SLIDE 11: POWER BI DASHBOARD =================
    slide = prs.slides.add_slide(blank_layout)
    add_header(slide, "Power BI Dashboard Overview")
    
    # Text details on left
    txBox_l = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_l = txBox_l.text_frame
    tf_l.word_wrap = True
    
    p_title = tf_l.paragraphs[0]
    p_title.text = "EXECUTIVE REPORTING HUB"
    p_title.font.name = 'Arial'
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = PRIMARY_BLUE
    p_title.space_after = Pt(15)
    
    bullets = [
        "<b>Page 1: Industry Overview:</b> high-level aggregation of net inflows, active folio growth, and treemap of category AUM.",
        "<b>Page 2: Fund Performance:</b> Return vs Volatility scatter plot, active scorecard ranking, NAV historical trends.",
        "<b>Page 3: Investor Analytics:</b> Heatmap of transaction amounts by state, age distribution vs ticket size.",
        "<b>Page 4: SIP & Market Trends:</b> Dual axis chart (SIP inflow vs Nifty index close value) showing market momentum correlations."
    ]
    for b in bullets:
        p_b = tf_l.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(13)
        p_b.font.color.rgb = TEXT_DARK
        p_b.space_after = Pt(12)
        
    # Right column: custom visual graphic or placeholder showing Dashboard Design details
    txBox_r = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.9), Inches(4.8))
    tf_r = txBox_r.text_frame
    tf_r.word_wrap = True
    
    p_title_r = tf_r.paragraphs[0]
    p_title_r.text = "DASHBOARD DESIGN SYSTEM"
    p_title_r.font.name = 'Arial'
    p_title_r.font.size = Pt(18)
    p_title_r.font.bold = True
    p_title_r.font.color.rgb = SECONDARY_ORANGE
    p_title_r.space_after = Pt(15)
    
    # Styled block for design spec
    spec_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.5), Inches(5.9), Inches(3.6))
    spec_box.fill.solid()
    spec_box.fill.fore_color.rgb = BG_LIGHT
    spec_box.line.color.rgb = PRIMARY_BLUE
    tf_spec = spec_box.text_frame
    tf_spec.word_wrap = True
    p_s = tf_spec.paragraphs[0]
    p_s.text = "UI SPECIFICATION (bluestock_theme.json)\n\n• Primary Color: #414BEA (Bluestock Blue)\n• Secondary Color: #F05537 (Bluestock Coral/Orange)\n• Gridlines: Light Gray (#eeeeee)\n• Border Corner Radius: 10px rounded card borders\n• Font Family: Segoe UI for uniform modern styling"
    p_s.font.name = 'Arial'
    p_s.font.size = Pt(14)
    p_s.font.color.rgb = TEXT_DARK

    # ================= SLIDE 12: CONCLUSION & RECOMMENDATIONS (Dark Background) =================
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, PRIMARY_BLUE)
    
    # Large colored vertical accent block
    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.3), Inches(7.5))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = SECONDARY_ORANGE
    accent_bar.line.fill.background()
    
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = "CONCLUSION & FUTURE SCOPE"
    p.font.name = 'Arial'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.space_after = Pt(20)
    
    bullets = [
        "<b>1. Machine Learning recommendations:</b> integrate a collaborative filtering model into the investor portal to recommend mutual funds based on demographic profiles.",
        "<b>2. Churn Prevention Model:</b> train a Random Forest classifier using computed SIP continuity day gaps to flag users likely to stop systematic payments.",
        "<b>3. Portfolio Optimization:</b> add a modern portfolio theory (MPT) optimizer allowing users to adjust asset weights to maximize Sharpe ratio."
    ]
    for b in bullets:
        p_b = tf.add_paragraph()
        p_b.text = b
        p_b.font.name = 'Arial'
        p_b.font.size = Pt(16)
        p_b.font.color.rgb = RGBColor(0xee, 0xee, 0xee)
        p_b.space_after = Pt(15)
        
    p_ty = tf.add_paragraph()
    p_ty.text = "\nTHANK YOU"
    p_ty.font.name = 'Arial'
    p_ty.font.size = Pt(24)
    p_ty.font.bold = True
    p_ty.font.color.rgb = SECONDARY_ORANGE
    
    # Save Presentation
    prs.save(output_pptx_path)
    print("PowerPoint presentation generated successfully!")
    
    # Copy copy
    try:
        import shutil
        shutil.copy(output_pptx_path, os.path.join(reports_dir, 'Bluestock_MF_Presentation.pptx'))
        print("Copied PPTX presentation to reports/")
    except Exception as e:
        print(f"Failed to copy PPTX presentation: {e}")

if __name__ == '__main__':
    create_presentation()
